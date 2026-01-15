from typing import List
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt
from pptx.util import Inches

from core.models.slide_spec import SlideSpec
from core.models.theme import SlideTheme


def _add_body_box(slide):
    return slide.shapes.add_textbox(
        Inches(1.0),
        Inches(1.8),
        Inches(8.0),
        Inches(4.5),
    )


def render_presentation(
    slides: List[SlideSpec],
    theme: SlideTheme,
    output_path: Path,
) -> Path:
    """
    Render a list of SlideSpec objects into a PowerPoint file
    using a predefined theme.
    """

    prs = Presentation()

    for slide_spec in slides:
        layout_idx = theme.layouts.get(slide_spec.slide_type, 1)

        if slide_spec.slide_type == "text":
            _add_text_slide(prs, slide_spec, layout_idx, theme)

        elif slide_spec.slide_type == "chart":
            _add_chart_slide(prs, slide_spec, layout_idx, theme)

        elif slide_spec.slide_type == "table":
            _add_table_slide(prs, slide_spec, layout_idx, theme)

    prs.save(output_path)
    return output_path


def _apply_title_style(shape, theme: SlideTheme):
    p = shape.text_frame.paragraphs[0]
    p.font.name = theme.fonts["title"]["name"]
    p.font.size = Pt(theme.fonts["title"]["size"])


def _apply_body_style(shape, theme: SlideTheme):
    for p in shape.text_frame.paragraphs:
        p.font.name = theme.fonts["body"]["name"]
        p.font.size = Pt(theme.fonts["body"]["size"])


def _add_text_slide(
    prs: Presentation,
    spec: SlideSpec,
    layout_idx: int,
    theme: SlideTheme,
):
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])

    title = slide.shapes.title
    body = _add_body_box(slide)

    title.text = spec.title
    _apply_title_style(title, theme)

    config = spec.config or {}

    lines = config.get("text")
    if not lines:
        # fallback to description
        lines = [spec.description] if spec.description else []

    body.text = "\n".join(lines)
    _apply_body_style(body, theme)


def _add_chart_slide(
    prs: Presentation,
    spec: SlideSpec,
    layout_idx: int,
    theme: SlideTheme,
):
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])

    title = slide.shapes.title
    body = _add_body_box(slide)

    title.text = spec.title
    _apply_title_style(title, theme)

    # Placeholder: textual chart representation
    lines = [
        " | ".join(f"{k}: {v}" for k, v in row.items())
        for row in spec.config.get("data", [])
    ]

    body.text = "\n".join(lines)
    _apply_body_style(body, theme)


def _add_table_slide(
    prs: Presentation,
    spec: SlideSpec,
    layout_idx: int,
    theme: SlideTheme,
):
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])

    title = slide.shapes.title
    body = _add_body_box(slide)

    title.text = spec.title
    _apply_title_style(title, theme)

    body.text = str(spec.config)
    _apply_body_style(body, theme)
